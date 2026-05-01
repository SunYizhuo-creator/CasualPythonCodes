# ========== 修复1：替换xlrd为openpyxl（支持xlsx格式） ==========
import openpyxl
import os
from pptx import Presentation

# Excel文件路径（使用原始字符串避免转义）
xlsx_path = r"D:\python练手示例\小精灵\information.xlsx"

# 读取Excel工作簿（改用openpyxl）
wb = openpyxl.load_workbook(xlsx_path)
# 获取第一个工作表
ws = wb.worksheets[0]

# ========== 修复2：修正图片文件夹路径（使用Windows格式） ==========
# 图片文件夹路径（改为Windows格式的原始字符串）
pathFile = r"D:\python练手示例\小精灵\picture"
# 获取路径下所有文件名
allItems = os.listdir(pathFile)

# PPT文件路径
pptx_path = r"D:\python练手示例\小精灵\简历.pptx"
# 读取PPT
pptxFile = Presentation(pptx_path)
# 获取母版中第一个版式
layoutOne = pptxFile.slide_layouts[0]
# 获取所有幻灯片
slideAll = pptxFile.slides

# 先新建一个页面，确认版式中各个占位符的序号及类型（调试用，可保留）
newSlide = slideAll.add_slide(layoutOne)
for single_plh in newSlide.placeholders:
    phf = single_plh.placeholder_format
    print(f"占位符序号：{phf.idx}，类型：{phf.type}")

# ========== 修复3：修正Excel行遍历逻辑（openpyxl的行索引从1开始） ==========
# 遍历第2~6行（对应Excel的第2行到第6行，索引1到5）
for row in range(2, 7):  # openpyxl中行号从1开始，range(2,7)对应2-6行
    # 获取当前行的单元格值
    name = ws.cell(row=row, column=1).value  # 第1列：姓名
    major = ws.cell(row=row, column=2).value  # 第2列：专业
    mobile = ws.cell(row=row, column=3).value  # 第3列：手机号
    school = ws.cell(row=row, column=4).value  # 第4列：学校
    comment = ws.cell(row=row, column=5).value  # 第5列：备注
    
    # 处理空值，避免写入PPT时出错
    name = name if name else ""
    major = major if major else ""
    mobile = str(mobile) if mobile else ""
    school = school if school else ""
    comment = comment if comment else ""
    
    # 遍历图片文件，匹配姓名
    for item in allItems:
        # 分割文件名和后缀
        file_name = item.split(".")[0]
        # 如果文件名与姓名匹配
        if file_name == name:
            # 新增幻灯片
            slide = slideAll.add_slide(layoutOne) 
            
            # ========== 注意：请根据实际占位符序号修改以下数字 ==========
            # 填入文本内容（序号请替换为你实际查到的占位符序号）
            slide.placeholders[13].text = name
            slide.placeholders[15].text = major
            slide.placeholders[16].text = mobile
            slide.placeholders[17].text = school
            slide.placeholders[18].text = comment
            
            # 拼接图片路径（修复路径拼接逻辑）
            picPath = os.path.join(pathFile, item)  # 更安全的路径拼接
            # 检查图片文件是否存在
            if os.path.exists(picPath):
                # 在指定占位符中插入图片（序号请确认是否正确）
                slide.placeholders[1].insert_picture(picPath)
            else:
                print(f"警告：图片文件不存在 - {picPath}")

# 保存PPT文件
pptxFile.save(pptx_path)
print("成功结束")