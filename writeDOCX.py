# 使用from...import从pptx模块中导入Presentation
from pptx import Presentation
# 使用import导入docx
import docx


# 必须导入：qn函数（设置中文字体）、Pt（设置字号）
from docx.shared import Pt
from docx.oxml.ns import qn  # 核心修复：添加qn的导入


# 新建一个空白Word文档，赋值给变量docxFile
docxFile = docx.Document()


normal_style = docxFile.styles['Normal']
# 西文字体（兼容英文/数字）
normal_style.font.name = '微软雅黑'
# 中文字体（关键：解决中文不生效问题）
normal_style._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
# 可选：设置默认字号，避免字体大小混乱
normal_style.font.size = Pt(11)



# 将.pptx文件路径赋值给变量path
path = r"C:\Users\Administrator\OneDrive\高中学习经验\64级实际适用课件之1-强基（孙逸卓)  .pptx"
# 读取path并赋值给变量pptxFile
pptxFile = Presentation(path)

# 将变量n设置为1
n = 1

# for循环遍历pptxFile中的.slides属性，赋值给slide
for slide in pptxFile.slides:
    
    # 向文档中添加标题f"第{n}页"，为二级标题
    docxFile.add_heading(f"第{n}页",level=2)

    # for循环遍历slide中.shapes属性，赋值给变量shape
    for shape in slide.shapes:
        # 判断形状中是否有文本框
        if shape.has_text_frame == True:
            # 读取形状中的文本框，并赋值给变量textFrame
            textFrame = shape.text_frame
        
            # for循环遍历文本框内的所有段落
            # 赋值给变量paragraph
            for paragraph in textFrame.paragraphs:
                # for循环遍历段落中的所有样式块
                # 赋值给变量run
                for run in paragraph.runs:
                    # 读取样式块中的文本内容
                    texts = run.text
                    # 向Word文档中添加段落texts的文本内容
                    docxFile.add_paragraph(texts)

    # 将变量n进行累加
    n = n + 1

# 保存文档到指定路径，并命名为"资料.docx"
docxFile.save(r"C:\Users\Administrator\OneDrive\编程成果\DOCX.docx")
print("成功运行完成")